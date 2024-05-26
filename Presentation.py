from pptx import Presentation

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Threat Actor Analysis: OCEANLOTUS"
subtitle.text = "Cybersecurity Course Final Project\nYour Name\nDate"

# Slide 2: Classification of OCEANLOTUS
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Classification of OCEANLOTUS"
content.text = "- Aliases: APT32, Cobalt Kitty, SeaLotus, APT-C-00\n- Origin: Vietnam\n- Skill Level: High\n- Resources: Significant, likely state-sponsored\n- Classification: Advanced Persistent Threat (APT)"

# Slide 3: Motivations of OCEANLOTUS
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Motivations of OCEANLOTUS"
content.text = "- Political Espionage: Targeting government entities, dissidents, and foreign corporations\n- Economic Espionage: Focusing on sectors such as manufacturing, consumer products, and hospitality\n- Geo-Political Context: Operates within a framework of state interests, likely linked to Vietnamese strategic objectives"

# Slide 4: Geo-Political Insight
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Geo-Political Insight"
content.text = "- Regional Focus: Southeast Asia, including China and neighboring countries\n- Strategic Goals: Gathering intelligence on political opponents and foreign entities to maintain state security and economic advantage"

# Slide 5: Tradecraft and Tactics
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Tradecraft and Tactics"
content.text = "- Reconnaissance: Spear-phishing, social engineering\n- Weaponization: Customized malware, exploit kits\n- Delivery: Phishing emails, malicious attachments\n- Exploitation: Exploiting vulnerabilities in software and systems\n- Installation: Remote access tools (RATs), backdoors\n- Command and Control: Use of compromised servers and domains for communication\n- Actions on Objectives: Data exfiltration, system manipulation"

# Slide 6: The Lockheed Martin Kill Chain
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "The Lockheed Martin Kill Chain"
content.text = "- Reconnaissance: Identifying targets through open-source intelligence (OSINT)\n- Weaponization: Developing malware tailored to specific vulnerabilities\n- Delivery: Sending phishing emails with malicious links or attachments\n- Exploitation: Exploiting software vulnerabilities to gain initial access\n- Installation: Installing malware to establish a foothold\n- Command and Control: Establishing communication with compromised systems\n- Actions on Objectives: Exfiltrating data, maintaining persistent access"

# Slide 7: Case Study: 2017 Attack on a Global Corporation
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Case Study: 2017 Attack on a Global Corporation"
content.text = "- Primary Effect: Unauthorized access to sensitive data\n- Secondary Effect: Financial losses due to data breach\n- Second Order Effect: Damage to the company's reputation and loss of customer trust"

# Slide 8: Case Study: 2019 Attack on Government Agencies
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Case Study: 2019 Attack on Government Agencies"
content.text = "- Primary Effect: Compromise of government communication systems\n- Secondary Effect: Disruption of governmental operations\n- Second Order Effect: Potential national security implications"

# Slide 9: Primary Effects of OCEANLOTUS Attacks
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Primary Effects of OCEANLOTUS Attacks"
content.text = "- Data Breach: Unauthorized access to sensitive information\n- System Compromise: Disruption of IT systems and services"

# Slide 10: Secondary Effects of OCEANLOTUS Attacks
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Secondary Effects of OCEANLOTUS Attacks"
content.text = "- Financial Losses: Costs associated with remediation and recovery\n- Operational Disruption: Impact on business continuity and service delivery"

# Slide 11: Second Order Effects of OCEANLOTUS Attacks
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Second Order Effects of OCEANLOTUS Attacks"
content.text = "- Reputation Damage: Long-term harm to organizational credibility\n- Economic Impact: Broader economic consequences due to disrupted supply chains"

# Slide 12: Strategic Concerns for Policy Makers
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Strategic Concerns for Policy Makers"
content.text = "- Public Concern: OCEANLOTUS poses a significant threat to national security and economic stability\n- Private Problem: Businesses face direct financial and operational risks"

# Slide 13: Policy Recommendations
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Policy Recommendations"
content.text = "- Strengthening Cyber Defenses: Investment in cybersecurity infrastructure and training\n- International Collaboration: Working with global partners to address the threat\n- Regulation and Compliance: Implementing stringent cybersecurity regulations and compliance requirements"

# Slide 14: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = "- Summary: OCEANLOTUS is a highly skilled and resourceful APT with significant impacts on both public and private sectors.\n- Call to Action: Enhanced vigilance, robust defense mechanisms, and international cooperation are essential to mitigate the threat."

# Slide 15: References
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "References"
content.text = "- Sources: FireEye, Symantec, CrowdStrike, various cybersecurity reports and publications\n- Citations: APA or MLA format as required"

# Save the presentation
prs.save("OCEANLOTUS_Threat_Actor_Analysis.pptx")
